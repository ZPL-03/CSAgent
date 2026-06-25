from __future__ import annotations

import json
import gzip
import sys

from core.knowledge_ingestion import (
    KnowledgeIngestionService,
    SUPPORTED_INGEST_SUFFIXES,
    SUPPORTED_QT_FILE_FILTER,
    _mineru_backend_and_method_for,
    _run_external_parser,
)


def test_runtime_knowledge_empty_status_exposes_pipeline_contract(tmp_path) -> None:
    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=64, chunk_overlap_tokens=12)

    status = service.status()

    assert status["ready"] is True
    assert status["builtin_ready"] is True
    assert status["store_type"] == "project_runtime_knowledge"
    assert status["document_count"] == 0
    assert status["runtime_document_count"] == 0
    assert status["runtime_rag_chunk_count"] == 0
    assert status["runtime_kg_entity_count"] == 0
    assert status["runtime_kg_relation_count"] == 0
    assert status["builtin_rag_chunk_count"] > 0
    assert status["builtin_kg_entity_count"] > 0
    assert status["builtin_kg_relation_count"] > 0
    assert status["rag_chunk_count"] == status["builtin_rag_chunk_count"]
    assert status["vector_chunk_count"] == 0
    assert status["vector_ready"] is False
    assert status["vector_status"] == "pending"
    assert status["vector_message"] == "等待写入向量索引"
    assert status["vector_detail"] == ""
    assert status["vector_backend"] == ""
    assert status["kg_entity_count"] == status["builtin_kg_entity_count"]
    assert status["kg_relation_count"] == status["builtin_kg_relation_count"]
    assert status["chunk_token_size"] == 64
    assert status["chunk_overlap_tokens"] == 12
    assert status["dedupe_key"] == "content_hash"
    assert [step["status"] for step in status["pipeline"]] == ["pending", "pending", "pending", "pending", "pending"]
    assert [step["name"] for step in status["pipeline"]] == [
        "MinerU / Docling 文档解析",
        "语义分块",
        "BGE-M3 向量化索引",
        "KG 实体/关系抽取",
        "检索验证 / 证据引用",
    ]


def test_runtime_knowledge_ingestion_builds_chunks_kg_vector_index_and_dedupes(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    source = tmp_path / "pressure_hull_notes.md"
    source.write_text(
        "\n\n".join(
            [
                "# 复合材料外压圆柱耐压壳设计",
                "T700 碳纤维/环氧复合材料 pressure hull 在 external pressure 下需要关注 buckling 和 imperfection sensitivity。",
                "ASME RD-1172 用于线性屈曲压力估算，PBIPF 用于极限压力预测，Abaqus Lanczos 与 Riks 用于有限元验证。",
                "制造上需要控制 filament winding 张力、固化过程和铺层角度偏差，避免分层和初始缺陷放大。",
                "这些信息用于 RAG 检索、KG 实体关系抽取和报告解释，不替代 FEM 结果。",
            ]
        ),
        encoding="utf-8",
    )

    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=48, chunk_overlap_tokens=8)
    first = service.ingest_file(source)
    second = service.ingest_file(source)

    assert first.success
    assert second.success
    assert first.document_id == second.document_id
    assert first.chunk_count >= 1
    assert first.entity_count >= 3
    assert first.relation_count >= 1
    assert first.retrieval_verification["status"] == "success"
    assert first.retrieval_verification["evidence_chunks"]

    manifest = json.loads(service.manifest_path.read_text(encoding="utf-8"))
    assert manifest["document_count"] == 1
    assert manifest["rag_chunk_count"] == first.chunk_count
    assert manifest["chunk_token_size"] == 48
    assert manifest["chunk_overlap_tokens"] == 8
    assert manifest["dedupe_key"] == "content_hash"
    assert manifest["vector_chunk_count"] == first.chunk_count
    assert manifest["vector_collection_name"] == "csdm_cph_project_knowledge"
    assert manifest["last_ingestion"]["steps"][2]["status"] == "success"
    assert manifest["last_ingestion"]["steps"][4]["name"] == "检索验证 / 证据引用"
    assert manifest["last_ingestion"]["steps"][4]["status"] == "success"
    assert manifest["last_ingestion"]["retrieval_verification"]["hit_count"] >= 1
    assert manifest["last_retrieval_verification"]["evidence_chunks"]

    chunks = [json.loads(line) for line in service.chunks_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    assert len(chunks) == first.chunk_count
    assert all(item["token_estimate"] <= 56 for item in chunks)
    assert len({item["content_hash"] for item in chunks}) == len(chunks)

    relations = [json.loads(line) for line in service.relations_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    chunk_ids = {item["chunk_id"] for item in chunks}
    assert all(item["evidence_chunk_id"] in chunk_ids for item in relations)


def test_runtime_knowledge_ingestion_prefers_mineru_and_falls_back_to_docling(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")

    first_source = tmp_path / "mineru_priority.pdf"
    first_source.write_bytes(b"%PDF-1.4\n% temporary parser routing fixture\n")

    def mineru_parser(self, path):
        return "\n\n".join(
            [
                "# MinerU pressure hull",
                "T700 composite pressure hull uses ASME RD-1172, PBIPF and Abaqus Riks.",
                "Filament winding and curing quality control reduce initial imperfection risk.",
            ]
        )

    def docling_must_not_run(self, path):
        raise AssertionError("Docling should not run when MinerU succeeds.")

    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_mineru", mineru_parser)
    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_docling", docling_must_not_run)
    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge_mineru", chunk_token_size=48, chunk_overlap_tokens=8)
    first_result = service.ingest_file(first_source)

    assert first_result.success
    assert first_result.parser_backend == "mineru"
    assert first_result.chunk_count >= 1

    second_source = tmp_path / "docling_fallback.pdf"
    second_source.write_bytes(b"%PDF-1.4\n% temporary parser routing fixture\n")

    def mineru_unavailable(self, path):
        raise RuntimeError("mineru unavailable")

    def docling_parser(self, path):
        return "\n\n".join(
            [
                "# Docling pressure hull",
                "M40J composite pressure hull under external pressure uses ASME RD-1172 and PBIPF.",
                "Abaqus finite element verification records buckling and ultimate pressure evidence.",
            ]
        )

    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_mineru", mineru_unavailable)
    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_docling", docling_parser)
    fallback_service = KnowledgeIngestionService(
        base_dir=tmp_path / "knowledge_docling",
        chunk_token_size=48,
        chunk_overlap_tokens=8,
    )
    second_result = fallback_service.ingest_file(second_source)

    assert second_result.success
    assert second_result.parser_backend == "docling"
    assert second_result.chunk_count >= 1


def test_external_parser_bypasses_local_api_proxy(monkeypatch) -> None:
    monkeypatch.delenv("NO_PROXY", raising=False)
    monkeypatch.delenv("no_proxy", raising=False)

    returncode, stdout, _stderr = _run_external_parser(
        [sys.executable, "-c", "import os; print(os.environ.get('NO_PROXY', ''))"],
        15,
    )

    assert returncode == 0
    values = {item.strip() for item in stdout.strip().split(",")}
    assert {"127.0.0.1", "localhost", "::1"}.issubset(values)


def test_mineru_text_pdf_uses_pipeline_text_route(tmp_path, monkeypatch) -> None:
    monkeypatch.delenv("MINERU_BACKEND", raising=False)
    monkeypatch.delenv("MINERU_PARSE_METHOD", raising=False)
    from reportlab.pdfgen import canvas

    source = tmp_path / "text_layer.pdf"
    pdf = canvas.Canvas(str(source))
    for row in range(20):
        pdf.drawString(72, 760 - row * 24, "CSAgent pressure hull text layer. " * 5)
    pdf.save()

    backend, method = _mineru_backend_and_method_for(source)

    assert backend == "pipeline"
    assert method == "txt"


def test_mineru_sparse_pdf_uses_high_accuracy_route(tmp_path, monkeypatch) -> None:
    monkeypatch.delenv("MINERU_BACKEND", raising=False)
    monkeypatch.delenv("MINERU_PARSE_METHOD", raising=False)
    from reportlab.pdfgen import canvas

    source = tmp_path / "sparse_layer.pdf"
    pdf = canvas.Canvas(str(source))
    pdf.showPage()
    pdf.save()

    backend, method = _mineru_backend_and_method_for(source)

    assert backend == "hybrid-auto-engine"
    assert method == "auto"


def test_runtime_knowledge_ingestion_accepts_office_table_and_pdf_text_layers(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")

    def external_parser_unavailable(self, path):
        raise RuntimeError("external parser disabled for test")

    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_mineru", external_parser_unavailable)
    monkeypatch.setattr(KnowledgeIngestionService, "_parse_with_docling", external_parser_unavailable)

    import docx
    import openpyxl
    from pptx import Presentation
    from reportlab.pdfgen import canvas

    docx_path = tmp_path / "pressure_hull_report.docx"
    document = docx.Document()
    document.add_paragraph("T700 composite pressure hull uses ASME RD-1172, PBIPF and Abaqus Riks verification.")
    document.add_paragraph("Filament winding and curing quality control reduce imperfection and delamination risk.")
    document.save(str(docx_path))

    pptx_path = tmp_path / "pressure_hull_review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Composite pressure hull review"
    text_box = slide.shapes.add_textbox(914400, 1371600, 7315200, 914400)
    text_box.text = "Buckling, PBIPF and finite element verification evidence."
    presentation.save(str(pptx_path))

    xlsx_path = tmp_path / "pressure_hull_table.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "design"
    sheet.append(["material", "formula", "verification"])
    sheet.append(["T800G", "ASME RD-1172", "Abaqus"])
    workbook.save(xlsx_path)

    pdf_path = tmp_path / "pressure_hull_notes.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 720, "Composite pressure hull PDF text layer: ASME RD-1172, PBIPF, Abaqus Riks.")
    pdf.drawString(72, 700, "Manufacturing evidence includes filament winding and curing control.")
    pdf.save()

    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=48, chunk_overlap_tokens=8)
    results = {
        path.suffix: service.ingest_file(path)
        for path in [docx_path, pptx_path, xlsx_path, pdf_path]
    }

    assert results[".docx"].success
    assert results[".docx"].parser_backend == "python_docx_text"
    assert results[".pptx"].success
    assert results[".pptx"].parser_backend == "python_pptx_text"
    assert results[".xlsx"].success
    assert results[".xlsx"].parser_backend == "openpyxl_table"
    assert results[".pdf"].success
    assert results[".pdf"].parser_backend == "pdf_text"

    manifest = json.loads(service.manifest_path.read_text(encoding="utf-8"))
    assert manifest["document_count"] == 4
    assert manifest["rag_chunk_count"] >= 4
    assert manifest["kg_relation_count"] >= 1
    assert manifest["last_retrieval_verification"]["evidence_chunks"]


def test_runtime_knowledge_ingestion_marks_vector_status_warning_when_backend_fails(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")

    class BrokenRAGEngine:
        def __init__(self, *args, **kwargs) -> None:
            raise RuntimeError("vector backend unavailable")

    source = tmp_path / "pressure_hull_vector_failure.md"
    source.write_text(
        "\n\n".join(
            [
                "# 外压耐压壳向量状态",
                "T700 composite pressure hull under external pressure uses ASME RD-1172 for buckling checks.",
                "PBIPF and Abaqus Riks provide ultimate pressure verification.",
            ]
        ),
        encoding="utf-8",
    )

    monkeypatch.setattr(
        "core.knowledge_ingestion._build_rag_engine",
        lambda **kwargs: BrokenRAGEngine(**kwargs),
    )
    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=48, chunk_overlap_tokens=8)
    result = service.ingest_file(source)

    manifest = json.loads(service.manifest_path.read_text(encoding="utf-8"))
    assert result.success
    assert manifest["vector_status"] == "warning"
    assert manifest["vector_ready"] is False
    assert manifest["vector_chunk_count"] == 0
    assert manifest["vector_backend"] == "unavailable"
    assert manifest["vector_message"] == "向量索引写入失败"


def test_runtime_knowledge_ingestion_emits_step_progress(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    source = tmp_path / "pressure_hull_progress.md"
    source.write_text(
        "\n\n".join(
            [
                "# 复合材料外压圆柱耐压壳",
                "T800G composite pressure hull 在 external pressure 下需要关注 buckling 和 imperfection sensitivity。",
                "ASME RD-1172、PBIPF、Abaqus Lanczos 与 Static Riks 用于校核设计结果。",
                "filament winding、curing 和铺层角度偏差是制造质量控制重点。",
            ]
        ),
        encoding="utf-8",
    )

    progress_events: list[list[dict]] = []

    service = KnowledgeIngestionService(
        base_dir=tmp_path / "knowledge",
        chunk_token_size=40,
        chunk_overlap_tokens=6,
        progress_callback=lambda steps: progress_events.append([dict(step) for step in steps]),
    )
    result = service.ingest_file(source)

    assert result.success
    assert len(progress_events) >= 4
    assert progress_events[0][0]["status"] == "running"
    assert any(event[1]["status"] == "running" for event in progress_events)
    assert any(event[2]["status"] == "running" for event in progress_events)
    assert progress_events[-1][0]["status"] == "success"
    assert progress_events[-1][1]["status"] == "success"
    assert progress_events[-1][2]["status"] == "success"
    assert progress_events[-1][3]["status"] in {"success", "warning"}
    assert progress_events[-1][4]["name"] == "检索验证 / 证据引用"
    assert progress_events[-1][4]["status"] == "success"


def test_runtime_knowledge_rebuild_and_snapshot_export(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    source = tmp_path / "pressure_hull_rebuild.md"
    source.write_text(
        "\n\n".join(
            [
                "# 外压耐压壳知识入库",
                "composite pressure hull 在 external pressure 下需要结合 ASME RD-1172、PBIPF 与 Abaqus 校核。",
                "buckling、postbuckling、initial imperfection 与 filament winding 质量控制都需要进入证据链。",
            ]
        ),
        encoding="utf-8",
    )

    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=42, chunk_overlap_tokens=6)
    result = service.ingest_file(source)
    assert result.success

    rebuilt = service.rebuild_indexes()

    assert rebuilt["document_count"] == 1
    assert rebuilt["rag_chunk_count"] >= 1
    assert rebuilt["kg_entity_count"] >= 1
    assert rebuilt["kg_relation_count"] >= 1
    assert rebuilt["last_reindex"]["duplicate_chunk_count"] == 0
    assert rebuilt["last_retrieval_verification"]["hit_count"] >= 1
    assert rebuilt["pipeline"][0]["name"] == "MinerU / Docling 文档解析"
    assert rebuilt["pipeline"][2]["name"] == "BGE-M3 向量化索引"

    snapshot_path = service.export_snapshot()
    payload = json.loads(snapshot_path.read_text(encoding="utf-8"))

    assert payload["schema"] == "csagent_project_knowledge_snapshot_v1"
    assert payload["manifest"]["document_count"] == 1
    assert len(payload["documents"]) == 1
    assert len(payload["chunks"]) == rebuilt["rag_chunk_count"]
    assert len(payload["entities"]) == rebuilt["kg_entity_count"]
    assert len(payload["relations"]) == rebuilt["kg_relation_count"]


def test_runtime_knowledge_ingestion_merges_new_document_with_existing_rag_and_kg(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    first = tmp_path / "asme_pressure_hull.md"
    first.write_text(
        "\n\n".join(
            [
                "# ASME pressure hull",
                "T700 composite pressure hull under external pressure uses ASME RD-1172 for buckling checks.",
                "PBIPF and Abaqus Riks provide ultimate pressure verification.",
            ]
        ),
        encoding="utf-8",
    )
    second = tmp_path / "manufacturing_pressure_hull.md"
    second.write_text(
        "\n\n".join(
            [
                "# manufacturing pressure hull",
                "Filament winding, curing and fiber placement quality control reduce imperfection and delamination risk.",
                "Abaqus finite element verification links manufacturing defects with buckling and postbuckling behavior.",
            ]
        ),
        encoding="utf-8",
    )

    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=42, chunk_overlap_tokens=6)
    first_result = service.ingest_file(first)
    second_result = service.ingest_file(second)

    assert first_result.success
    assert second_result.success
    assert first_result.document_id != second_result.document_id

    documents = [json.loads(line) for line in service.documents_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    chunks = [json.loads(line) for line in service.chunks_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    entities = [json.loads(line) for line in service.entities_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    relations = [json.loads(line) for line in service.relations_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    manifest = json.loads(service.manifest_path.read_text(encoding="utf-8"))

    assert {item["document_id"] for item in documents} == {first_result.document_id, second_result.document_id}
    assert {item["source_id"] for item in chunks} == {first_result.document_id, second_result.document_id}
    assert any(first_result.document_id in item.get("sources", []) for item in entities)
    assert any(second_result.document_id in item.get("sources", []) for item in entities)
    assert any(item.get("record_id") == first_result.document_id for item in relations)
    assert any(item.get("record_id") == second_result.document_id for item in relations)
    assert manifest["document_count"] == 2
    assert manifest["rag_chunk_count"] == len(chunks)
    assert manifest["kg_entity_count"] == len(entities)
    assert manifest["kg_relation_count"] == len(relations)


def test_runtime_knowledge_status_reports_builtin_runtime_and_merged_counts(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    builtin_dir = tmp_path / "csllm"
    builtin_rag_path = builtin_dir / "rag" / "rag_chunks.compact.jsonl.gz"
    builtin_kg_dir = builtin_dir / "kg"
    builtin_manifest_path = builtin_dir / "provenance" / "manifest.json"
    builtin_rag_path.parent.mkdir(parents=True, exist_ok=True)
    with gzip.open(builtin_rag_path, "wt", encoding="utf-8") as handle:
        handle.write(
            json.dumps(
                {
                    "chunk_id": "BASE_CHUNK_1",
                    "record_id": "BASE_DOC_1",
                    "source_id": "BASE_DOC_1",
                    "retrieval_scope": "main",
                    "document_title": "Built pressure hull reference",
                    "content_plain": "composite pressure hull ASME RD-1172 buckling",
                },
                ensure_ascii=False,
            )
            + "\n"
        )
    builtin_kg_dir.mkdir(parents=True, exist_ok=True)
    (builtin_kg_dir / "entities.jsonl").write_text(
        json.dumps({"type": "DesignFormula", "name": "ASME RD-1172"}, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    with gzip.open(builtin_kg_dir / "relations.compact.jsonl.gz", "wt", encoding="utf-8") as handle:
        handle.write(
            json.dumps(
                {
                    "source_type": "DesignFormula",
                    "source": "ASME RD-1172",
                    "target_type": "FailureMode",
                    "target": "Buckling",
                    "relation": "PREDICTED_BY",
                    "record_id": "BASE_DOC_1",
                },
                ensure_ascii=False,
            )
            + "\n"
        )
    (builtin_kg_dir / "kg_stats.json").write_text(
        json.dumps({"total_entities": 1, "total_relations": 1}, ensure_ascii=False),
        encoding="utf-8",
    )
    builtin_manifest_path.parent.mkdir(parents=True, exist_ok=True)
    builtin_manifest_path.write_text(
        json.dumps({"record_counts": {"rag_chunks": 1, "entities": 1, "relations": 1}}, ensure_ascii=False),
        encoding="utf-8",
    )

    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=42, chunk_overlap_tokens=6)
    service.builtin_dir = builtin_dir
    service.builtin_rag_chunks_path = builtin_rag_path
    service.builtin_kg_dir = builtin_kg_dir
    service.builtin_manifest_path = builtin_manifest_path

    source = tmp_path / "uploaded_pressure_hull.md"
    source.write_text(
        "# uploaded pressure hull\n\nT700 composite pressure hull under external pressure uses PBIPF and Abaqus Riks.",
        encoding="utf-8",
    )
    result = service.ingest_file(source)
    status = service.status()
    snapshot_path = service.export_snapshot()
    snapshot = json.loads(snapshot_path.read_text(encoding="utf-8"))

    assert result.success
    assert status["ready"] is True
    assert status["builtin_ready"] is True
    assert status["builtin_rag_chunk_count"] == 1
    assert status["runtime_document_count"] == 1
    assert status["runtime_rag_chunk_count"] == result.chunk_count
    assert status["rag_chunk_count"] == result.chunk_count + 1
    assert status["kg_relation_count"] == result.relation_count + 1
    assert len(snapshot["documents"]) == 1
    assert len(snapshot["chunks"]) == result.chunk_count
    assert all(chunk.get("source_id") == result.document_id for chunk in snapshot["chunks"])


def test_runtime_knowledge_ingestion_accepts_engineering_text_and_binary_metadata(tmp_path, monkeypatch) -> None:
    monkeypatch.setenv("CSDM_cph_USE_HASH_EMBEDDING", "1")
    service = KnowledgeIngestionService(base_dir=tmp_path / "knowledge", chunk_token_size=48, chunk_overlap_tokens=8)

    status_file = tmp_path / "buckling_job.sta"
    status_file.write_text(
        "Abaqus Lanczos buckling step completed for composite pressure hull. PBIPF and ASME RD-1172 are checked.",
        encoding="utf-8",
    )
    status_result = service.ingest_file(status_file)

    odb_file = tmp_path / "buckling_job.odb"
    odb_file.write_bytes(b"ODB_BINARY_PLACEHOLDER")
    odb_result = service.ingest_file(odb_file)

    assert status_result.success
    assert status_result.parser_backend == "text"
    assert odb_result.success
    assert odb_result.parser_backend == "engineering_metadata"
    assert odb_result.chunk_count >= 1

    chunks = [json.loads(line) for line in service.chunks_path.read_text(encoding="utf-8").splitlines() if line.strip()]
    odb_chunks = [item for item in chunks if item.get("source_id") == odb_result.document_id]
    assert odb_chunks
    assert any("工程二进制文件元数据" in item.get("content_markdown", "") for item in odb_chunks)


def test_runtime_knowledge_supported_suffix_contract_is_shared_with_gui_filter() -> None:
    for suffix in [".pdf", ".docx", ".pptx", ".xlsx", ".png", ".webp", ".inp", ".sta", ".odb"]:
        assert suffix in SUPPORTED_INGEST_SUFFIXES
        assert f"*{suffix}" in SUPPORTED_QT_FILE_FILTER
