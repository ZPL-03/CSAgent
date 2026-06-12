"""项目知识库入库流水线。"""

from __future__ import annotations

import csv
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

from core.config_loader import load_app_config
from core.domain_knowledge import trim_text
from core.paths import (
    CHROMA_DIR,
    KNOWLEDGE_UPLOADS_DIR,
    RUNTIME_KNOWLEDGE_DIR,
    RUNTIME_KNOWLEDGE_KG_DIR,
    RUNTIME_KNOWLEDGE_RAG_DIR,
    RUNTIME_KNOWLEDGE_STRUCTURED_DIR,
)
from core.rag_engine import RAGEngine


TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".inp",
    ".py",
    ".for",
    ".f",
    ".f90",
    ".log",
    ".jnl",
    ".rpy",
    ".dat",
    ".msg",
    ".sta",
    ".prt",
}
TABLE_SUFFIXES = {".csv", ".tsv"}
PDF_SUFFIXES = {".pdf"}
DOCX_SUFFIXES = {".docx"}
PPTX_SUFFIXES = {".pptx"}
EXCEL_SUFFIXES = {".xlsx", ".xlsm"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
BINARY_ENGINEERING_SUFFIXES = {".odb", ".sim", ".cae"}
MINERU_SUFFIXES = PDF_SUFFIXES | DOCX_SUFFIXES | PPTX_SUFFIXES | IMAGE_SUFFIXES
DOCLING_SUFFIXES = PDF_SUFFIXES | DOCX_SUFFIXES | PPTX_SUFFIXES | IMAGE_SUFFIXES
SUPPORTED_INGEST_SUFFIXES = (
    TEXT_SUFFIXES
    | TABLE_SUFFIXES
    | PDF_SUFFIXES
    | DOCX_SUFFIXES
    | PPTX_SUFFIXES
    | EXCEL_SUFFIXES
    | IMAGE_SUFFIXES
    | BINARY_ENGINEERING_SUFFIXES
)
SUPPORTED_QT_FILE_FILTER = "知识资料 (" + " ".join(f"*{suffix}" for suffix in sorted(SUPPORTED_INGEST_SUFFIXES)) + ");;所有文件 (*.*)"

STEP_PARSE = "MinerU / Docling 文档解析"
STEP_CHUNK = "语义分块"
STEP_VECTOR = "BGE-M3 向量化索引"
STEP_KG = "Neo4j 实体/关系抽取"
STEP_RETRIEVAL = "检索验证 / 证据引用"
_TIKTOKEN_ENCODING: Any | None = None
_TIKTOKEN_READY = False

DOMAIN_ENTITY_TERMS: dict[str, dict[str, str]] = {
    "Material": {
        "T700": "T700",
        "T800": "T800",
        "T800G": "T800G",
        "M40J": "M40J",
        "碳纤维": "Carbon Fiber",
        "环氧": "Epoxy",
        "CFRP": "CFRP",
        "composite": "Composite Material",
        "carbon fiber": "Carbon Fiber",
    },
    "Structure": {
        "耐压壳": "Pressure Hull",
        "压力壳": "Pressure Hull",
        "圆柱壳": "Cylindrical Shell",
        "pressure hull": "Pressure Hull",
        "cylindrical shell": "Cylindrical Shell",
        "laminate": "Laminate",
        "铺层": "Layup",
    },
    "FailureMode": {
        "屈曲": "Buckling",
        "后屈曲": "Postbuckling",
        "collapse": "Collapse",
        "buckling": "Buckling",
        "delamination": "Delamination",
        "分层": "Delamination",
        "缺陷": "Initial Imperfection",
        "imperfection": "Initial Imperfection",
    },
    "DesignFormula": {
        "ASME": "ASME RD-1172",
        "RD-1172": "ASME RD-1172",
        "PBIPF": "PBIPF",
        "NASA": "NASA SP-8007",
    },
    "VerificationMethod": {
        "Abaqus": "Abaqus",
        "有限元": "Finite Element Analysis",
        "finite element": "Finite Element Analysis",
        "Riks": "Static Riks",
        "Lanczos": "Lanczos Buckling",
    },
    "ManufacturingProcess": {
        "缠绕": "Filament Winding",
        "铺放": "Fiber Placement",
        "固化": "Curing",
        "filament winding": "Filament Winding",
        "fiber placement": "Fiber Placement",
        "curing": "Curing",
    },
}


@dataclass
class PipelineStep:
    name: str
    status: str = "pending"
    message: str = ""
    detail: str = ""


@dataclass
class IngestionResult:
    document_id: str
    title: str
    source_path: str
    stored_path: str
    parser_backend: str
    markdown_path: str
    chunk_count: int
    entity_count: int
    relation_count: int
    duplicate_chunk_count: int
    retrieval_verification: dict[str, Any]
    steps: list[PipelineStep]

    @property
    def success(self) -> bool:
        return all(step.status in {"success", "warning"} for step in self.steps)


def _utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _safe_stem(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9_\-\u4e00-\u9fff]+", "_", value).strip("_")
    return text[:80] or "document"


def _read_text(path: Path) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _write_jsonl(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as handle:
        for row in rows:
            handle.write(json.dumps(row, ensure_ascii=False) + "\n")


def _read_jsonl(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    rows: list[dict[str, Any]] = []
    with path.open("r", encoding="utf-8") as handle:
        for line in handle:
            if not line.strip():
                continue
            try:
                payload = json.loads(line)
            except json.JSONDecodeError:
                continue
            if isinstance(payload, dict):
                rows.append(payload)
    return rows


def _markdown_plain_text(markdown: str) -> str:
    text = re.sub(r"!\[[^\]]*\]\([^)]+\)", " ", markdown)
    text = re.sub(r"```.*?```", " ", text, flags=re.DOTALL)
    text = re.sub(r"^#+\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"[$*_`>#|]+", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _token_count(text: str) -> int:
    """按 tiktoken 优先估算 token 数；缺失时使用中英文混合近似值。"""
    global _TIKTOKEN_ENCODING, _TIKTOKEN_READY
    if not _TIKTOKEN_READY:
        try:
            import tiktoken
        except ModuleNotFoundError:
            _TIKTOKEN_ENCODING = None
        else:
            try:
                _TIKTOKEN_ENCODING = tiktoken.get_encoding("cl100k_base")
            except Exception:
                _TIKTOKEN_ENCODING = None
        _TIKTOKEN_READY = True
    if _TIKTOKEN_ENCODING is not None:
        return max(1, len(_TIKTOKEN_ENCODING.encode(text)))
    chinese_chars = len(re.findall(r"[\u4e00-\u9fff]", text))
    other_chars = max(0, len(text) - chinese_chars)
    return max(1, chinese_chars + round(other_chars / 4))


def _split_text_by_tokens(text: str, max_tokens: int) -> list[str]:
    """把单个超长文本块按句子和 token 预算切开。"""
    if _token_count(text) <= max_tokens:
        return [text]
    separators = r"(?<=[。！？；;.!?])\s+|(?<=\.)\s+(?=[A-Z])|\n+"
    sentences = [part.strip() for part in re.split(separators, text) if part.strip()]
    if len(sentences) <= 1:
        sentences = [part.strip() for part in re.split(r"(?<=，)|(?<=,)\s*", text) if part.strip()]
    parts: list[str] = []
    current: list[str] = []
    for sentence in sentences:
        probe = "\n".join([*current, sentence]).strip()
        if current and _token_count(probe) > max_tokens:
            parts.append("\n".join(current).strip())
            current = [sentence]
        elif _token_count(sentence) > max_tokens:
            if current:
                parts.append("\n".join(current).strip())
                current = []
            parts.extend(_split_long_string(sentence, max_tokens))
        else:
            current.append(sentence)
    if current:
        parts.append("\n".join(current).strip())
    return [part for part in parts if part.strip()]


def _split_long_string(text: str, max_tokens: int) -> list[str]:
    parts: list[str] = []
    current: list[str] = []
    for char in text:
        probe = "".join([*current, char])
        if current and _token_count(probe) > max_tokens:
            parts.append("".join(current).strip())
            current = [char]
        else:
            current.append(char)
    if current:
        parts.append("".join(current).strip())
    return [part for part in parts if part]


class KnowledgeIngestionService:
    """把用户上传资料转为本项目可检索的 RAG/KG 运行时数据。"""

    def __init__(
        self,
        base_dir: Path | None = None,
        *,
        chunk_token_size: int | None = None,
        chunk_overlap_tokens: int | None = None,
        min_chunk_tokens: int | None = None,
        progress_callback: Callable[[list[dict[str, Any]]], None] | None = None,
    ) -> None:
        app_config = load_app_config()
        knowledge_config = dict(app_config.get("project_knowledge", {}))
        configured_base = Path(str(knowledge_config.get("base_dir", RUNTIME_KNOWLEDGE_DIR)))
        if not configured_base.is_absolute():
            configured_base = RUNTIME_KNOWLEDGE_DIR.parent.parent / configured_base
        self.base_dir = base_dir or configured_base
        self.uploads_dir = self.base_dir / "uploads"
        self.structured_dir = self.base_dir / "structured_text"
        self.markdown_dir = self.structured_dir / "markdown_documents"
        self.rag_dir = self.base_dir / "rag"
        self.kg_dir = self.base_dir / "kg"
        self.manifest_path = self.base_dir / "manifest.json"
        self.documents_path = self.structured_dir / "documents.jsonl"
        self.blocks_path = self.structured_dir / "blocks.jsonl"
        self.chunks_path = self.rag_dir / "rag_chunks.jsonl"
        self.entities_path = self.kg_dir / "entities.jsonl"
        self.relations_path = self.kg_dir / "relations.jsonl"
        self.stats_path = self.kg_dir / "kg_stats.json"
        self.chunk_token_size = int(chunk_token_size or knowledge_config.get("chunk_token_size", 512) or 512)
        self.chunk_overlap_tokens = int(chunk_overlap_tokens or knowledge_config.get("chunk_overlap_tokens", 64) or 64)
        self.min_chunk_tokens = int(min_chunk_tokens or knowledge_config.get("min_chunk_tokens", 80) or 80)
        self.vector_enabled = bool(knowledge_config.get("vector_enabled", True))
        self.vector_collection_name = str(knowledge_config.get("vector_collection_name", "csdm_cph_project_knowledge"))
        configured_chroma_dir = knowledge_config.get("vector_chroma_dir")
        if configured_chroma_dir:
            chroma_dir = Path(str(configured_chroma_dir))
            self.vector_chroma_dir = chroma_dir if chroma_dir.is_absolute() else RUNTIME_KNOWLEDGE_DIR.parent.parent / chroma_dir
        elif base_dir is not None:
            self.vector_chroma_dir = self.base_dir / "chroma_db"
        else:
            self.vector_chroma_dir = CHROMA_DIR
        self.progress_callback = progress_callback

    def _notify_progress(self, steps: list[PipelineStep]) -> None:
        if self.progress_callback is not None:
            self.progress_callback([asdict(step) for step in steps])

    def ensure_dirs(self) -> None:
        for path in [
            self.base_dir,
            self.uploads_dir,
            self.structured_dir,
            self.markdown_dir,
            self.rag_dir,
            self.kg_dir,
            KNOWLEDGE_UPLOADS_DIR,
            RUNTIME_KNOWLEDGE_STRUCTURED_DIR,
            RUNTIME_KNOWLEDGE_RAG_DIR,
            RUNTIME_KNOWLEDGE_KG_DIR,
        ]:
            path.mkdir(parents=True, exist_ok=True)

    def ingest_file(self, source_path: str | Path) -> IngestionResult:
        self.ensure_dirs()
        source = Path(source_path).expanduser().resolve()
        if not source.exists() or not source.is_file():
            raise FileNotFoundError(f"资料文件不存在：{source}")

        file_hash = _sha256(source)
        document_id = f"DOC_{file_hash[:12]}"
        stored_name = f"{document_id}_{_safe_stem(source.stem)}{source.suffix.lower()}"
        stored_path = self.uploads_dir / stored_name
        if source != stored_path:
            shutil.copy2(source, stored_path)

        steps = [
            PipelineStep(STEP_PARSE, "running", "正在解析资料"),
            PipelineStep(STEP_CHUNK),
            PipelineStep(STEP_VECTOR),
            PipelineStep(STEP_KG),
            PipelineStep(STEP_RETRIEVAL),
        ]
        self._notify_progress(steps)

        try:
            markdown, parser_backend = self._parse_document(stored_path)
        except Exception as exc:
            steps[0] = PipelineStep(STEP_PARSE, "failed", "解析失败", str(exc))
            self._notify_progress(steps)
            result = IngestionResult(
                document_id=document_id,
                title=source.stem,
                source_path=str(source),
                stored_path=str(stored_path),
                parser_backend="",
                markdown_path="",
                chunk_count=0,
                entity_count=0,
                relation_count=0,
                duplicate_chunk_count=0,
                retrieval_verification={},
                steps=steps,
            )
            self._write_manifest(last_result=result)
            raise

        markdown = self._clean_markdown(markdown)
        if not markdown:
            steps[0] = PipelineStep(STEP_PARSE, "failed", "解析结果为空", "解析器没有返回可检索文本。")
            self._notify_progress(steps)
            result = IngestionResult(
                document_id=document_id,
                title=source.stem,
                source_path=str(source),
                stored_path=str(stored_path),
                parser_backend=parser_backend,
                markdown_path="",
                chunk_count=0,
                entity_count=0,
                relation_count=0,
                duplicate_chunk_count=0,
                retrieval_verification={},
                steps=steps,
            )
            self._write_manifest(last_result=result)
            raise RuntimeError(f"解析结果为空：{source.name}")

        steps[0] = PipelineStep(STEP_PARSE, "success", f"{parser_backend} 解析完成", f"Markdown 长度 {len(markdown)} 字符")
        steps[1] = PipelineStep(
            STEP_CHUNK,
            "running",
            "正在生成语义文本块",
            f"{self.chunk_token_size} token 目标窗口，overlap={self.chunk_overlap_tokens} token",
        )
        self._notify_progress(steps)

        markdown_path = self.markdown_dir / f"{document_id}.md"
        markdown_path.write_text(markdown, encoding="utf-8")
        blocks = self._build_blocks(document_id, source.stem, stored_path, markdown, parser_backend)
        chunks = self._build_chunks(document_id, source.stem, stored_path, blocks)
        entities, relations = self._build_kg(document_id, source.stem, chunks)
        steps[1] = PipelineStep(
            STEP_CHUNK,
            "success",
            f"生成 {len(chunks)} 个候选 RAG 文本块",
            f"{len(blocks)} 个结构化文本块",
        )
        steps[2] = PipelineStep(STEP_VECTOR, "running", "正在写入向量索引", self.vector_collection_name)
        steps[3] = PipelineStep(STEP_KG, "running", "正在整理实体/关系", f"候选实体 {len(entities)} 个，候选关系 {len(relations)} 条")
        steps[4] = PipelineStep(STEP_RETRIEVAL, "pending", "等待索引和关系写入")
        self._notify_progress(steps)

        write_stats = self._replace_document_records(
            document_id,
            source,
            stored_path,
            parser_backend,
            markdown_path,
            blocks,
            chunks,
            entities,
            relations,
        )

        steps[0] = PipelineStep(STEP_PARSE, "success", f"{parser_backend} 解析完成", f"Markdown 长度 {len(markdown)} 字符")
        steps[1] = PipelineStep(
            STEP_CHUNK,
            "success",
            f"生成 {write_stats['chunk_count']} 个 RAG 文本块",
            f"{self.chunk_token_size} token 目标窗口，overlap={self.chunk_overlap_tokens} token，去重 {write_stats['duplicate_chunk_count']} 个",
        )
        steps[2] = PipelineStep(
            STEP_VECTOR,
            write_stats["vector_status"],
            write_stats["vector_message"],
            write_stats["vector_detail"],
        )
        steps[3] = PipelineStep(
            STEP_KG,
            "success" if write_stats["relation_count"] else "warning",
            f"抽取 {write_stats['entity_count']} 个实体、{write_stats['relation_count']} 条关系",
            "规则词典抽取；写入本地实体和关系 JSONL",
        )
        retrieval_verification = dict(write_stats.get("retrieval_verification") or {})
        steps[4] = PipelineStep(
            STEP_RETRIEVAL,
            str(write_stats.get("retrieval_status") or "warning"),
            str(write_stats.get("retrieval_message") or "检索验证状态未知"),
            str(write_stats.get("retrieval_detail") or ""),
        )
        self._notify_progress(steps)
        result = IngestionResult(
            document_id=document_id,
            title=source.stem,
            source_path=str(source),
            stored_path=str(stored_path),
            parser_backend=parser_backend,
            markdown_path=str(markdown_path),
            chunk_count=write_stats["chunk_count"],
            entity_count=write_stats["entity_count"],
            relation_count=write_stats["relation_count"],
            duplicate_chunk_count=write_stats["duplicate_chunk_count"],
            retrieval_verification=retrieval_verification,
            steps=steps,
        )
        self._write_manifest(last_result=result)
        return result

    def status(self) -> dict[str, Any]:
        manifest = self._load_manifest()
        return {
            "base_dir": str(self.base_dir),
            "upload_dir": str(self.uploads_dir),
            "manifest_path": str(self.manifest_path),
            "rag_chunks_path": str(self.chunks_path),
            "kg_dir": str(self.kg_dir),
            "ready": self.chunks_path.exists() and int(manifest.get("rag_chunk_count", 0) or 0) > 0,
            **manifest,
        }

    def rebuild_indexes(self) -> dict[str, Any]:
        """基于当前已解析文本块重建向量索引、KG 关系、统计清单和检索验证状态。"""

        self.ensure_dirs()
        documents = _read_jsonl(self.documents_path)
        chunks, duplicate_chunk_count = self._dedupe_chunks(_read_jsonl(self.chunks_path))
        retained_chunk_ids = {str(row.get("chunk_id") or "") for row in chunks}
        if len(chunks) != len(_read_jsonl(self.chunks_path)):
            _write_jsonl(self.chunks_path, chunks)

        entities: list[dict[str, Any]] = []
        relations: list[dict[str, Any]] = []
        document_titles = {str(row.get("document_id") or ""): str(row.get("title") or row.get("file_name") or row.get("document_id") or "") for row in documents}
        for document_id, title in document_titles.items():
            document_chunks = [row for row in chunks if str(row.get("source_id") or "") == document_id]
            doc_entities, doc_relations = self._build_kg(document_id, title, document_chunks)
            entities.extend(doc_entities)
            relations.extend(doc_relations)
        entities = self._dedupe_entities(entities)
        relations = self._dedupe_relations(
            [row for row in relations if str(row.get("evidence_chunk_id") or "") in retained_chunk_ids]
        )
        _write_jsonl(self.entities_path, entities)
        _write_jsonl(self.relations_path, relations)

        vector_stats = self._sync_vector_index(chunks)
        verification = (
            self._verify_retrieval_evidence(
                str(documents[-1].get("document_id") or ""),
                str(documents[-1].get("title") or documents[-1].get("file_name") or ""),
                chunks,
                relations,
            )
            if documents
            else {
                "status": "warning",
                "message": "当前没有可验证的入库资料",
                "detail": "document_count=0",
                "hit_count": 0,
                "relation_hit_count": 0,
                "invalid_relation_refs": [],
                "evidence_chunks": [],
                "evidence_relations": [],
                "verified_at": _utc_now(),
            }
        )
        entity_counts: dict[str, int] = {}
        relation_counts: dict[str, int] = {}
        for entity in entities:
            entity_type = str(entity.get("type") or "Unknown")
            entity_counts[entity_type] = entity_counts.get(entity_type, 0) + 1
        for relation in relations:
            relation_type = str(relation.get("relation") or "UNKNOWN")
            relation_counts[relation_type] = relation_counts.get(relation_type, 0) + 1
        self.stats_path.write_text(
            json.dumps(
                {
                    "total_entities": len(entities),
                    "total_relations": len(relations),
                    "entity_counts": entity_counts,
                    "relation_type_counts": relation_counts,
                    "updated_at": _utc_now(),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        steps = [
            PipelineStep(STEP_PARSE, "success" if documents else "warning", f"复用 {len(documents)} 份已解析资料"),
            PipelineStep(STEP_CHUNK, "success" if chunks else "warning", f"复用 {len(chunks)} 个去重文本块", f"去重 {duplicate_chunk_count} 个重复文本块"),
            PipelineStep(
                STEP_VECTOR,
                str(vector_stats.get("status") or "warning"),
                str(vector_stats.get("message") or "向量索引状态未知"),
                str(vector_stats.get("detail") or ""),
            ),
            PipelineStep(STEP_KG, "success" if relations else "warning", f"重建 {len(entities)} 个实体、{len(relations)} 条关系"),
            PipelineStep(
                STEP_RETRIEVAL,
                str(verification.get("status") or "warning"),
                str(verification.get("message") or ""),
                str(verification.get("detail") or ""),
            ),
        ]
        self._write_manifest()
        manifest = self._load_manifest()
        manifest["pipeline"] = [asdict(step) for step in steps]
        manifest["last_reindex"] = {
            "updated_at": _utc_now(),
            "document_count": len(documents),
            "rag_chunk_count": len(chunks),
            "duplicate_chunk_count": duplicate_chunk_count,
            "kg_entity_count": len(entities),
            "kg_relation_count": len(relations),
            "vector_status": vector_stats,
            "retrieval_verification": verification,
        }
        manifest["last_retrieval_verification"] = verification
        self.manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
        return manifest

    def export_snapshot(self, output_path: str | Path | None = None) -> Path:
        """导出当前项目知识库快照，包含文档、文本块、实体、关系、统计和 manifest。"""

        self.ensure_dirs()
        if output_path is None:
            output_dir = self.base_dir / "snapshots"
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = output_dir / f"knowledge_snapshot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        target = Path(output_path)
        if not target.is_absolute():
            target = self.base_dir / target
        target.parent.mkdir(parents=True, exist_ok=True)
        payload = {
            "schema": "csagent_project_knowledge_snapshot_v1",
            "exported_at": _utc_now(),
            "manifest": self._load_manifest(),
            "documents": _read_jsonl(self.documents_path),
            "blocks": _read_jsonl(self.blocks_path),
            "chunks": _read_jsonl(self.chunks_path),
            "entities": _read_jsonl(self.entities_path),
            "relations": _read_jsonl(self.relations_path),
            "kg_stats": self._load_json_file(self.stats_path),
        }
        target.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        return target

    def _load_json_file(self, path: Path) -> dict[str, Any]:
        if not path.exists():
            return {}
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return {}
        return payload if isinstance(payload, dict) else {}

    def _sync_vector_index(self, chunks: list[dict[str, Any]]) -> dict[str, Any]:
        if not self.vector_enabled:
            return {
                "status": "warning",
                "message": "向量索引未启用",
                "detail": "project_knowledge.vector_enabled=false；仍保留 JSONL/BM25 检索和KG关系。",
                "count": 0,
                "backend": "disabled",
            }
        documents: list[str] = []
        ids: list[str] = []
        metadatas: list[dict[str, Any]] = []
        for chunk in chunks:
            chunk_id = str(chunk.get("chunk_id") or "").strip()
            content = str(chunk.get("content_plain") or chunk.get("content_markdown") or chunk.get("text") or "").strip()
            if not chunk_id or not content:
                continue
            ids.append(chunk_id)
            documents.append(content)
            metadatas.append(
                {
                    "source": "PROJECT_KNOWLEDGE",
                    "chunk_id": chunk_id,
                    "record_id": chunk.get("record_id"),
                    "source_id": chunk.get("source_id"),
                    "document_title": chunk.get("document_title") or chunk.get("title"),
                    "title": chunk.get("title") or chunk.get("document_title"),
                    "section_title": chunk.get("section_title"),
                    "source_type": chunk.get("source_type"),
                    "parser_backend": chunk.get("parser_backend"),
                    "load_case_tag": chunk.get("load_case_tag"),
                    "design_platform_scope": chunk.get("design_platform_scope"),
                    "page_start": int(chunk.get("page_start") or 0),
                    "page_end": int(chunk.get("page_end") or 0),
                    "token_estimate": int(chunk.get("token_estimate") or 0),
                    "source_url": chunk.get("source_url") or "",
                    "doi": chunk.get("doi") or "",
                }
            )
        try:
            engine = RAGEngine(chroma_dir=self.vector_chroma_dir, collection_name=self.vector_collection_name)
            engine.reset_collection()
            engine.upsert_documents(ids=ids, documents=documents, metadatas=metadatas)
            backend = "hash_embedding" if engine.use_hash_embedding_only or engine._embedder_failed else engine.embedding_model_name
            return {
                "status": "success",
                "message": f"向量索引写入 {len(ids)} 个文本块",
                "detail": f"collection={self.vector_collection_name}，backend={backend}，目录={self.vector_chroma_dir}",
                "count": len(ids),
                "backend": backend,
            }
        except Exception as exc:
            return {
                "status": "warning",
                "message": "向量索引写入失败",
                "detail": trim_text(str(exc), 800),
                "count": 0,
                "backend": "unavailable",
            }

    def _load_manifest(self) -> dict[str, Any]:
        if not self.manifest_path.exists():
            return self._base_manifest()
        try:
            payload = json.loads(self.manifest_path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return self._base_manifest()
        if not isinstance(payload, dict):
            return self._base_manifest()
        return {**self._base_manifest(), **payload}

    def _parse_document(self, path: Path) -> tuple[str, str]:
        suffix = path.suffix.lower()
        errors: list[str] = []
        if suffix in MINERU_SUFFIXES:
            try:
                return self._parse_with_mineru(path), "mineru"
            except Exception as exc:
                errors.append(f"MinerU: {exc}")
        if suffix in DOCLING_SUFFIXES:
            try:
                return self._parse_with_docling(path), "docling"
            except Exception as exc:
                errors.append(f"Docling: {exc}")
        if suffix in PDF_SUFFIXES:
            try:
                return self._parse_pdf_text(path), "pypdf_text"
            except Exception as exc:
                errors.append(f"PDF 文本层: {exc}")
        if suffix in DOCX_SUFFIXES:
            try:
                return self._parse_docx_text(path), "python_docx_text"
            except Exception as exc:
                errors.append(f"DOCX 文本层: {exc}")
        if suffix in PPTX_SUFFIXES:
            try:
                return self._parse_pptx_text(path), "python_pptx_text"
            except Exception as exc:
                errors.append(f"PPTX 文本层: {exc}")
        if suffix in IMAGE_SUFFIXES:
            return self._parse_image_metadata(path, errors), "image_metadata"
        if suffix in BINARY_ENGINEERING_SUFFIXES:
            return self._parse_engineering_binary_metadata(path), "engineering_metadata"
        if suffix in TEXT_SUFFIXES:
            return _read_text(path), "text"
        if suffix in TABLE_SUFFIXES:
            return self._parse_csv_table(path), "table_text"
        if suffix in EXCEL_SUFFIXES:
            try:
                return self._parse_xlsx_table(path), "openpyxl_table"
            except Exception as exc:
                errors.append(f"表格解析: {exc}")
        raise RuntimeError("；".join(errors) if errors else f"暂不支持的资料格式：{suffix}")

    def _parse_with_mineru(self, path: Path) -> str:
        command = shutil.which("mineru")
        if command is None:
            executable_dir = Path(sys.executable).resolve().parent
            for candidate in [executable_dir / "mineru.exe", executable_dir / "Scripts" / "mineru.exe", executable_dir / "mineru"]:
                if candidate.exists():
                    command = str(candidate)
                    break
        if command is None:
            raise RuntimeError("当前环境没有可用的 MinerU 命令。")
        with tempfile.TemporaryDirectory(prefix="csagent_mineru_") as temp_name:
            output_dir = Path(temp_name)
            args = [
                command,
                "-p",
                str(path),
                "-o",
                str(output_dir),
                "-m",
                os.getenv("MINERU_PARSE_METHOD", "auto"),
                "-b",
                os.getenv("MINERU_BACKEND", "pipeline"),
                "--formula",
                "true",
                "--table",
                "true",
                "--image-analysis",
                "true",
            ]
            api_url = os.getenv("MINERU_API_URL", "").strip()
            if api_url:
                args.extend(["--api-url", api_url])
            completed = subprocess.run(
                args,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=int(os.getenv("CSAGENT_MINERU_TIMEOUT", "900")),
            )
            if completed.returncode != 0:
                message = "\n".join(part for part in [completed.stdout, completed.stderr] if part).strip()
                raise RuntimeError(trim_text(message, 1200) or "MinerU 命令返回非零状态。")
            markdown_files = sorted(output_dir.rglob("*.md"), key=lambda item: item.stat().st_size, reverse=True)
            markdown = "\n\n".join(item.read_text(encoding="utf-8", errors="replace") for item in markdown_files)
            if not markdown.strip():
                raise RuntimeError("MinerU 没有输出 Markdown。")
            return markdown

    def _parse_with_docling(self, path: Path) -> str:
        try:
            from docling.document_converter import DocumentConverter
            from docling.document_converter import ImageFormatOption, PdfFormatOption
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import PdfPipelineOptions
        except ModuleNotFoundError as exc:
            raise RuntimeError("当前环境没有安装 Docling。") from exc
        pdf_options = PdfPipelineOptions(do_ocr=True, do_table_structure=True, do_formula_enrichment=True, images_scale=2.0)
        converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pdf_options),
                InputFormat.IMAGE: ImageFormatOption(pipeline_options=pdf_options),
            }
        )
        result = converter.convert(str(path))
        document = getattr(result, "document", None)
        if document is None:
            raise RuntimeError("Docling 没有返回文档对象。")
        return str(document.export_to_markdown() or "")

    def _parse_pdf_text(self, path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ModuleNotFoundError as exc:
            raise RuntimeError("当前环境没有 pypdf，无法读取 PDF 文本层。") from exc
        reader = PdfReader(str(path))
        parts = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"## Page {index}\n\n{text.strip()}")
        if not parts:
            raise RuntimeError("PDF 文本层为空。")
        return "\n\n".join(parts)

    def _parse_docx_text(self, path: Path) -> str:
        try:
            import docx
        except ModuleNotFoundError as exc:
            raise RuntimeError("当前环境没有 python-docx，无法读取 DOCX 文本层。") from exc
        document = docx.Document(str(path))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table_index, table in enumerate(document.tables, start=1):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            markdown = self._rows_to_markdown(f"{path.stem}_table_{table_index}", rows)
            if markdown.strip():
                parts.append(markdown)
        if not parts:
            raise RuntimeError("DOCX 文本层为空。")
        return "\n\n".join(parts)

    def _parse_pptx_text(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ModuleNotFoundError as exc:
            raise RuntimeError("当前环境没有 python-pptx，无法读取 PPTX 文本层。") from exc
        presentation = Presentation(str(path))
        sections: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text or "").strip()
                    if text:
                        texts.append(text)
            if texts:
                sections.append(f"## Slide {slide_index}\n\n" + "\n\n".join(texts))
        if not sections:
            raise RuntimeError("PPTX 文本层为空。")
        return "\n\n".join(sections)

    def _parse_image_metadata(self, path: Path, parser_errors: list[str]) -> str:
        stat = path.stat()
        error_text = "\n".join(f"- {item}" for item in parser_errors) if parser_errors else "- 未调用外部 OCR 解析器。"
        return "\n\n".join(
            [
                f"# {path.stem}",
                "## 图片资料元数据",
                f"- 文件名：{path.name}",
                f"- 文件类型：{path.suffix.lower().lstrip('.')}",
                f"- 文件大小：{stat.st_size} bytes",
                "- 说明：MinerU / Docling 未返回可检索文本时，系统保留图片资料元数据，供资料追溯和后续重新解析。",
                "## 解析诊断",
                error_text,
            ]
        )

    def _parse_engineering_binary_metadata(self, path: Path) -> str:
        stat = path.stat()
        return "\n\n".join(
            [
                f"# {path.stem}",
                "## 工程二进制文件元数据",
                f"- 文件名：{path.name}",
                f"- 文件类型：{path.suffix.lower().lstrip('.')}",
                f"- 文件大小：{stat.st_size} bytes",
                "- 说明：该文件属于工程二进制资料，知识库保存可检索元数据；ODB/CAE 内部场变量和云图仍由有限元结果可视化链路读取。",
            ]
        )

    def _parse_csv_table(self, path: Path) -> str:
        dialect = "excel-tab" if path.suffix.lower() == ".tsv" else "excel"
        rows: list[list[str]] = []
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            reader = csv.reader(handle, dialect=dialect)
            for row in reader:
                rows.append([str(cell) for cell in row])
        return self._rows_to_markdown(path.stem, rows)

    def _parse_xlsx_table(self, path: Path) -> str:
        try:
            import openpyxl
        except ModuleNotFoundError as exc:
            raise RuntimeError("当前环境没有 openpyxl，无法解析 Excel。") from exc
        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell) for cell in row]
                if any(value.strip() for value in values):
                    rows.append(values)
            if rows:
                sections.append(self._rows_to_markdown(sheet.title, rows))
        return "\n\n".join(sections)

    def _rows_to_markdown(self, title: str, rows: list[list[str]]) -> str:
        if not rows:
            return ""
        width = max(len(row) for row in rows)
        padded = [row + [""] * (width - len(row)) for row in rows[:200]]
        header = padded[0]
        body = padded[1:] or [[""] * width]
        lines = [f"# {title}", "", "| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
        for row in body:
            lines.append("| " + " | ".join(cell.replace("\n", " ") for cell in row) + " |")
        return "\n".join(lines)

    def _clean_markdown(self, markdown: str) -> str:
        text = markdown.replace("\r\n", "\n").replace("\r", "\n")
        lines = [re.sub(r"[ \t]+", " ", line).rstrip() for line in text.splitlines()]
        text = "\n".join(lines)
        text = re.sub(r"\n{4,}", "\n\n\n", text)
        return text.strip()

    def _build_blocks(self, document_id: str, title: str, stored_path: Path, markdown: str, parser_backend: str) -> list[dict[str, Any]]:
        blocks: list[dict[str, Any]] = []
        order = 0
        current_heading = ""
        paragraphs = [part.strip() for part in re.split(r"\n\s*\n", markdown) if part.strip()]
        for paragraph in paragraphs:
            order += 1
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", paragraph)
            block_type = "heading" if heading_match else "text"
            if heading_match:
                current_heading = heading_match.group(2).strip()
            if paragraph.lstrip().startswith("|"):
                block_type = "table"
            block_id = f"{document_id}_B{order:04d}"
            blocks.append(
                {
                    "block_id": block_id,
                    "source_id": document_id,
                    "record_id": document_id,
                    "file_path": str(stored_path),
                    "file_name": stored_path.name,
                    "title": title,
                    "document_title": title,
                    "section_title": current_heading,
                    "section_path": [current_heading] if current_heading else [],
                    "block_type": block_type,
                    "parser_backend": parser_backend,
                    "order": order,
                    "text": paragraph,
                    "char_count": len(paragraph),
                    "parse_quality_score": self._quality_score(markdown),
                }
            )
        return blocks

    def _quality_score(self, markdown: str) -> float:
        plain = _markdown_plain_text(markdown)
        score = 55.0
        if len(plain) >= 1200:
            score += 20.0
        elif len(plain) >= 300:
            score += 12.0
        if re.search(r"^#{1,6}\s+", markdown, flags=re.MULTILINE):
            score += 8.0
        if "|" in markdown and "---" in markdown:
            score += 6.0
        if re.search(r"\\frac|\\alpha|\\beta|\$\$", markdown):
            score += 5.0
        return round(min(score, 96.0), 2)

    def _build_chunks(self, document_id: str, title: str, stored_path: Path, blocks: list[dict[str, Any]]) -> list[dict[str, Any]]:
        expanded_blocks: list[dict[str, Any]] = []
        for block in blocks:
            text = str(block.get("text") or "")
            parts = _split_text_by_tokens(text, self.chunk_token_size)
            if len(parts) == 1:
                expanded_blocks.append(block)
                continue
            for index, part in enumerate(parts, start=1):
                item = dict(block)
                item["text"] = part
                item["block_id"] = f"{block.get('block_id')}_P{index:03d}"
                item["char_count"] = len(part)
                expanded_blocks.append(item)

        chunks: list[dict[str, Any]] = []
        seen_hashes: set[str] = set()
        buffer: list[dict[str, Any]] = []
        buffer_tokens = 0

        def flush() -> None:
            nonlocal buffer, buffer_tokens
            if not buffer:
                return
            content = "\n\n".join(str(item.get("text") or "") for item in buffer).strip()
            if not content:
                buffer = []
                buffer_tokens = 0
                return
            index = len(chunks) + 1
            fingerprint = hashlib.sha256(content.encode("utf-8")).hexdigest()
            if fingerprint in seen_hashes:
                buffer = self._overlap_buffer(buffer)
                buffer_tokens = sum(_token_count(str(item.get("text") or "")) for item in buffer)
                return
            seen_hashes.add(fingerprint)
            section_path = buffer[-1].get("section_path") or []
            token_estimate = _token_count(content)
            chunks.append(
                {
                    "chunk_id": f"{document_id}_C{index:04d}_{fingerprint[:8]}",
                    "chunk_fingerprint": fingerprint,
                    "content_hash": fingerprint,
                    "record_id": document_id,
                    "source_id": document_id,
                    "chunk_type": "fulltext",
                    "retrieval_scope": "main",
                    "document_title": title,
                    "title": title,
                    "section_title": (section_path[-1] if section_path else ""),
                    "section_path": section_path,
                    "title_path": [title, *section_path] if section_path else [title],
                    "block_types": sorted({str(item.get("block_type") or "text") for item in buffer}),
                    "block_ids": [str(item.get("block_id")) for item in buffer],
                    "page_start": 0,
                    "page_end": 0,
                    "source_url": "",
                    "doi": "",
                    "source_type": stored_path.suffix.lower().lstrip(".") or "document",
                    "parser_backend": buffer[0].get("parser_backend"),
                    "parse_quality_score": buffer[0].get("parse_quality_score", 0),
                    "content_markdown": content,
                    "content_plain": _markdown_plain_text(content),
                    "text": trim_text(content, 1200),
                    "char_count": len(content),
                    "token_estimate": token_estimate,
                    "chunk_token_size": self.chunk_token_size,
                    "chunk_overlap_tokens": self.chunk_overlap_tokens,
                    "design_platform_scope": "pressure_hull_design",
                    "load_case_tag": "external_pressure",
                }
            )
            buffer = self._overlap_buffer(buffer)
            buffer_tokens = sum(_token_count(str(item.get("text") or "")) for item in buffer)

        for block in expanded_blocks:
            text = str(block.get("text") or "")
            token_len = _token_count(text)
            if buffer and buffer_tokens + token_len > self.chunk_token_size:
                flush()
            if buffer and buffer_tokens + token_len > self.chunk_token_size:
                buffer = []
                buffer_tokens = 0
            buffer.append(block)
            buffer_tokens += token_len
        flush()
        return self._merge_small_chunks(chunks)

    def _overlap_buffer(self, buffer: list[dict[str, Any]]) -> list[dict[str, Any]]:
        if self.chunk_overlap_tokens <= 0:
            return []
        overlap: list[dict[str, Any]] = []
        total = 0
        for item in reversed(buffer):
            token_len = _token_count(str(item.get("text") or ""))
            if token_len > self.chunk_overlap_tokens:
                break
            if overlap and total + token_len > self.chunk_overlap_tokens:
                break
            overlap.append(item)
            total += token_len
            if total >= self.chunk_overlap_tokens:
                break
        return list(reversed(overlap))

    def _merge_small_chunks(self, chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
        if self.min_chunk_tokens <= 0 or len(chunks) <= 1:
            return chunks
        merged: list[dict[str, Any]] = []
        for chunk in chunks:
            token_count = int(chunk.get("token_estimate", 0) or _token_count(str(chunk.get("content_markdown") or "")))
            if (
                merged
                and token_count < self.min_chunk_tokens
                and int(merged[-1].get("token_estimate", 0) or 0) + token_count <= self.chunk_token_size
            ):
                merged[-1] = self._combine_chunks(merged[-1], chunk)
                continue
            merged.append(chunk)
        return self._reindex_chunks(merged)

    def _combine_chunks(self, left: dict[str, Any], right: dict[str, Any]) -> dict[str, Any]:
        content = "\n\n".join(
            part.strip()
            for part in [str(left.get("content_markdown") or ""), str(right.get("content_markdown") or "")]
            if part.strip()
        )
        fingerprint = hashlib.sha256(content.encode("utf-8")).hexdigest()
        combined = dict(left)
        combined["chunk_fingerprint"] = fingerprint
        combined["content_hash"] = fingerprint
        combined["content_markdown"] = content
        combined["content_plain"] = _markdown_plain_text(content)
        combined["text"] = trim_text(content, 1200)
        combined["char_count"] = len(content)
        combined["token_estimate"] = _token_count(content)
        combined["block_ids"] = [*list(left.get("block_ids") or []), *list(right.get("block_ids") or [])]
        combined["block_types"] = sorted({*list(left.get("block_types") or []), *list(right.get("block_types") or [])})
        combined["section_path"] = right.get("section_path") or left.get("section_path") or []
        section_path = combined.get("section_path") or []
        combined["section_title"] = section_path[-1] if section_path else ""
        combined["title_path"] = [combined.get("title") or combined.get("document_title") or "", *section_path] if section_path else [combined.get("title") or combined.get("document_title") or ""]
        return combined

    def _reindex_chunks(self, chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
        result: list[dict[str, Any]] = []
        for index, chunk in enumerate(chunks, start=1):
            item = dict(chunk)
            fingerprint = str(item.get("content_hash") or item.get("chunk_fingerprint") or "")
            if not fingerprint:
                fingerprint = hashlib.sha256(str(item.get("content_markdown") or "").encode("utf-8")).hexdigest()
                item["chunk_fingerprint"] = fingerprint
                item["content_hash"] = fingerprint
            item["chunk_id"] = f"{item.get('record_id')}_C{index:04d}_{fingerprint[:8]}"
            result.append(item)
        return result

    def _build_kg(self, document_id: str, title: str, chunks: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
        entity_counts: dict[tuple[str, str], int] = {}
        entity_sources: dict[tuple[str, str], set[str]] = {}
        relations: list[dict[str, Any]] = []
        seen_relations: set[tuple[str, str, str, str]] = set()
        for chunk in chunks:
            text = f"{chunk.get('document_title')}\n{chunk.get('content_plain') or chunk.get('content_markdown')}"
            chunk_entities: list[tuple[str, str]] = []
            for entity_type, terms in DOMAIN_ENTITY_TERMS.items():
                for keyword, normalized in terms.items():
                    if keyword.lower() in text.lower():
                        key = (entity_type, normalized)
                        entity_counts[key] = entity_counts.get(key, 0) + 1
                        entity_sources.setdefault(key, set()).add(document_id)
                        if key not in chunk_entities:
                            chunk_entities.append(key)
            if document_id:
                lit_key = ("Literature", document_id)
                entity_counts[lit_key] = entity_counts.get(lit_key, 0) + 1
                entity_sources.setdefault(lit_key, set()).add(document_id)
            for entity_type, name in chunk_entities:
                rel_key = (name, "MENTIONED_IN", document_id, str(chunk.get("chunk_id")))
                if rel_key not in seen_relations:
                    seen_relations.add(rel_key)
                    relations.append(
                        {
                            "source_type": entity_type,
                            "source": name,
                            "target_type": "Literature",
                            "target": document_id,
                            "relation": "MENTIONED_IN",
                            "record_id": document_id,
                            "evidence_chunk_id": chunk.get("chunk_id"),
                            "evidence_chunk_fingerprint": chunk.get("chunk_fingerprint"),
                            "evidence_source_id": document_id,
                            "evidence_document_title": title,
                            "evidence_source_url": "",
                            "evidence_doi": "",
                            "page_start": 0,
                            "page_end": 0,
                            "section_path": chunk.get("section_path") or [],
                        }
                    )
            for left_index, left in enumerate(chunk_entities):
                for right in chunk_entities[left_index + 1 :]:
                    if left[1] == right[1]:
                        continue
                    rel_key = (left[1], "CO_OCCURS_WITH", right[1], str(chunk.get("chunk_id")))
                    if rel_key in seen_relations:
                        continue
                    seen_relations.add(rel_key)
                    relations.append(
                        {
                            "source_type": left[0],
                            "source": left[1],
                            "target_type": right[0],
                            "target": right[1],
                            "relation": "CO_OCCURS_WITH",
                            "record_id": document_id,
                            "evidence_chunk_id": chunk.get("chunk_id"),
                            "evidence_chunk_fingerprint": chunk.get("chunk_fingerprint"),
                            "evidence_source_id": document_id,
                            "evidence_document_title": title,
                            "evidence_source_url": "",
                            "evidence_doi": "",
                            "page_start": 0,
                            "page_end": 0,
                            "section_path": chunk.get("section_path") or [],
                        }
                    )
        entities = [
            {
                "id": f"{entity_type}:{name}",
                "type": entity_type,
                "name": name,
                "count": count,
                "sources": sorted(entity_sources.get((entity_type, name), set())),
            }
            for (entity_type, name), count in sorted(entity_counts.items(), key=lambda item: (item[0][0], item[0][1]))
        ]
        return entities, relations

    def _replace_document_records(
        self,
        document_id: str,
        source: Path,
        stored_path: Path,
        parser_backend: str,
        markdown_path: Path,
        blocks: list[dict[str, Any]],
        chunks: list[dict[str, Any]],
        entities: list[dict[str, Any]],
        relations: list[dict[str, Any]],
    ) -> dict[str, int]:
        document_record = {
            "document_id": document_id,
            "source_id": document_id,
            "record_id": document_id,
            "title": source.stem,
            "file_name": source.name,
            "source_path": str(source),
            "stored_path": str(stored_path),
            "markdown_path": str(markdown_path),
            "parser_backend": parser_backend,
            "file_sha256": _sha256(stored_path),
            "updated_at": _utc_now(),
            "chunk_count": len(chunks),
        }
        documents = [row for row in _read_jsonl(self.documents_path) if row.get("document_id") != document_id]
        documents.append(document_record)
        all_blocks = [row for row in _read_jsonl(self.blocks_path) if row.get("source_id") != document_id]
        all_blocks.extend(blocks)
        existing_chunks = [row for row in _read_jsonl(self.chunks_path) if row.get("source_id") != document_id]
        all_chunks, duplicate_chunk_count = self._dedupe_chunks([*existing_chunks, *chunks])
        retained_chunk_ids = {str(row.get("chunk_id") or "") for row in all_chunks}
        retained_new_chunk_ids = {str(row.get("chunk_id") or "") for row in chunks} & retained_chunk_ids
        all_entities = [row for row in _read_jsonl(self.entities_path) if document_id not in row.get("sources", [])]
        if retained_new_chunk_ids:
            all_entities.extend(entities)
        all_entities = self._dedupe_entities(all_entities)
        existing_relations = [row for row in _read_jsonl(self.relations_path) if row.get("record_id") != document_id]
        retained_new_relations = [
            row
            for row in relations
            if str(row.get("evidence_chunk_id") or "") in retained_new_chunk_ids
        ]
        all_relations = self._dedupe_relations([*existing_relations, *retained_new_relations])

        _write_jsonl(self.documents_path, documents)
        _write_jsonl(self.blocks_path, all_blocks)
        _write_jsonl(self.chunks_path, all_chunks)
        _write_jsonl(self.entities_path, all_entities)
        _write_jsonl(self.relations_path, all_relations)
        vector_stats = self._sync_vector_index(all_chunks)
        retrieval_verification = self._verify_retrieval_evidence(document_id, source.stem, all_chunks, all_relations)

        entity_counts: dict[str, int] = {}
        relation_counts: dict[str, int] = {}
        for entity in all_entities:
            entity_counts[str(entity.get("type") or "Unknown")] = entity_counts.get(str(entity.get("type") or "Unknown"), 0) + 1
        for relation in all_relations:
            relation_counts[str(relation.get("relation") or "UNKNOWN")] = relation_counts.get(str(relation.get("relation") or "UNKNOWN"), 0) + 1
        self.stats_path.write_text(
            json.dumps(
                {
                    "total_entities": len(all_entities),
                    "total_relations": len(all_relations),
                    "entity_counts": entity_counts,
                    "relation_type_counts": relation_counts,
                    "updated_at": _utc_now(),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        return {
            "chunk_count": len([row for row in all_chunks if row.get("source_id") == document_id]),
            "duplicate_chunk_count": duplicate_chunk_count,
            "entity_count": len([row for row in all_entities if document_id in row.get("sources", [])]),
            "relation_count": len([row for row in all_relations if row.get("record_id") == document_id]),
            "vector_count": int(vector_stats.get("count", 0) or 0),
            "vector_status": str(vector_stats.get("status") or "warning"),
            "vector_message": str(vector_stats.get("message") or "向量索引状态未知"),
            "vector_detail": str(vector_stats.get("detail") or ""),
            "vector_backend": str(vector_stats.get("backend") or ""),
            "retrieval_status": str(retrieval_verification.get("status") or "warning"),
            "retrieval_message": str(retrieval_verification.get("message") or ""),
            "retrieval_detail": str(retrieval_verification.get("detail") or ""),
            "retrieval_verification": retrieval_verification,
        }

    def _verify_retrieval_evidence(
        self,
        document_id: str,
        title: str,
        chunks: list[dict[str, Any]],
        relations: list[dict[str, Any]],
    ) -> dict[str, Any]:
        document_chunks = [row for row in chunks if row.get("source_id") == document_id]
        document_relations = [row for row in relations if row.get("record_id") == document_id]
        chunk_ids = {str(row.get("chunk_id") or "") for row in chunks}
        invalid_relations = [
            str(row.get("evidence_chunk_id") or "")
            for row in document_relations
            if str(row.get("evidence_chunk_id") or "") not in chunk_ids
        ]
        query_terms = self._verification_query_terms(title, document_chunks, document_relations)
        scored_chunks = self._score_chunks_for_verification(query_terms, chunks)
        top_hits = scored_chunks[:5]
        own_hits = [row for row in top_hits if row.get("source_id") == document_id]
        evidence_chunks = [
            {
                "chunk_id": row.get("chunk_id"),
                "source_id": row.get("source_id"),
                "document_title": row.get("document_title") or row.get("title"),
                "score": row.get("verification_score"),
                "text": trim_text(str(row.get("content_plain") or row.get("content_markdown") or row.get("text") or ""), 220),
            }
            for row in own_hits[:3]
        ]
        evidence_relations = [
            {
                "source": row.get("source"),
                "relation": row.get("relation"),
                "target": row.get("target"),
                "evidence_chunk_id": row.get("evidence_chunk_id"),
            }
            for row in document_relations[:3]
        ]
        status = "success"
        if not document_chunks or not own_hits or invalid_relations:
            status = "failed"
        elif not document_relations:
            status = "warning"
        message = (
            f"检索命中 {len(own_hits)} 个当前文档证据，引用 {len(document_relations)} 条图谱关系"
            if status != "failed"
            else "检索验证失败"
        )
        detail_parts = [
            f"query_terms={', '.join(query_terms[:8]) or '-'}",
            f"evidence_chunks={len(evidence_chunks)}",
            f"evidence_relations={len(evidence_relations)}",
        ]
        if invalid_relations:
            detail_parts.append(f"invalid_relation_refs={len(invalid_relations)}")
        return {
            "status": status,
            "message": message,
            "detail": "；".join(detail_parts),
            "query_terms": query_terms,
            "hit_count": len(own_hits),
            "relation_hit_count": len(document_relations),
            "invalid_relation_refs": invalid_relations,
            "evidence_chunks": evidence_chunks,
            "evidence_relations": evidence_relations,
            "verified_at": _utc_now(),
        }

    def _verification_query_terms(
        self,
        title: str,
        chunks: list[dict[str, Any]],
        relations: list[dict[str, Any]],
    ) -> list[str]:
        terms: list[str] = []
        for part in re.split(r"[\s_\-:：/\\|,，.。;；]+", title):
            item = part.strip()
            if len(item) >= 2:
                terms.append(item)
        for relation in relations[:12]:
            for key in ["source", "target"]:
                item = str(relation.get(key) or "").strip()
                if item and item not in terms:
                    terms.append(item)
        domain_terms = [
            "耐压壳",
            "外压",
            "圆柱壳",
            "pressure hull",
            "external pressure",
            "buckling",
            "PBIPF",
            "ASME",
            "Abaqus",
        ]
        text = "\n".join(str(row.get("content_plain") or row.get("content_markdown") or "") for row in chunks[:6]).lower()
        for term in domain_terms:
            if term.lower() in text and term not in terms:
                terms.append(term)
        return terms[:16]

    def _score_chunks_for_verification(self, query_terms: list[str], chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
        scored: list[dict[str, Any]] = []
        for chunk in chunks:
            content = " ".join(
                [
                    str(chunk.get("document_title") or chunk.get("title") or ""),
                    str(chunk.get("section_title") or ""),
                    str(chunk.get("content_plain") or chunk.get("content_markdown") or chunk.get("text") or ""),
                ]
            ).lower()
            score = 0.0
            for term in query_terms:
                normalized = term.lower()
                if normalized and normalized in content:
                    score += 1.0
            if score <= 0:
                continue
            item = dict(chunk)
            item["verification_score"] = round(score, 3)
            scored.append(item)
        return sorted(scored, key=lambda item: (float(item.get("verification_score") or 0.0), str(item.get("source_id") or "")), reverse=True)

    def _dedupe_chunks(self, chunks: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], int]:
        result: list[dict[str, Any]] = []
        seen: set[str] = set()
        duplicate_count = 0
        for chunk in chunks:
            content_hash = str(chunk.get("content_hash") or chunk.get("chunk_fingerprint") or "")
            if not content_hash:
                content_hash = hashlib.sha256(str(chunk.get("content_markdown") or chunk.get("text") or "").encode("utf-8")).hexdigest()
                chunk["content_hash"] = content_hash
            if content_hash in seen:
                duplicate_count += 1
                continue
            seen.add(content_hash)
            result.append(chunk)
        return result, duplicate_count

    def _dedupe_entities(self, entities: list[dict[str, Any]]) -> list[dict[str, Any]]:
        merged: dict[tuple[str, str], dict[str, Any]] = {}
        for entity in entities:
            key = (str(entity.get("type") or ""), str(entity.get("name") or ""))
            if not key[0] or not key[1]:
                continue
            current = merged.get(key)
            if current is None:
                item = dict(entity)
                item["sources"] = sorted({str(value) for value in item.get("sources", []) if str(value)})
                merged[key] = item
                continue
            current["count"] = int(current.get("count", 0) or 0) + int(entity.get("count", 0) or 0)
            sources = {str(value) for value in current.get("sources", []) if str(value)}
            sources.update(str(value) for value in entity.get("sources", []) if str(value))
            current["sources"] = sorted(sources)
        return sorted(merged.values(), key=lambda item: (str(item.get("type")), str(item.get("name"))))

    def _dedupe_relations(self, relations: list[dict[str, Any]]) -> list[dict[str, Any]]:
        result: list[dict[str, Any]] = []
        seen: set[tuple[str, str, str, str, str]] = set()
        for relation in relations:
            key = (
                str(relation.get("source") or ""),
                str(relation.get("relation") or ""),
                str(relation.get("target") or ""),
                str(relation.get("record_id") or ""),
                str(relation.get("evidence_chunk_id") or ""),
            )
            if key in seen:
                continue
            seen.add(key)
            result.append(relation)
        return result

    def _write_manifest(self, last_result: IngestionResult | None = None) -> None:
        documents = _read_jsonl(self.documents_path)
        chunks = _read_jsonl(self.chunks_path)
        entities = _read_jsonl(self.entities_path)
        relations = _read_jsonl(self.relations_path)
        blocks = _read_jsonl(self.blocks_path)
        manifest = {
            **self._base_manifest(),
            "document_count": len(documents),
            "rag_chunk_count": len(chunks),
            "vector_chunk_count": len(chunks) if self.vector_enabled else 0,
            "vector_ready": bool(self.vector_enabled and chunks),
            "kg_entity_count": len(entities),
            "kg_relation_count": len(relations),
            "structured_document_count": len(documents),
            "structured_block_count": len(blocks),
            "markdown_document_count": len(list(self.markdown_dir.glob("*.md"))) if self.markdown_dir.exists() else 0,
            "updated_at": _utc_now(),
        }
        if last_result is not None:
            manifest["last_ingestion"] = {
                "document_id": last_result.document_id,
                "title": last_result.title,
                "success": last_result.success,
                "parser_backend": last_result.parser_backend,
                "chunk_count": last_result.chunk_count,
                "entity_count": last_result.entity_count,
                "relation_count": last_result.relation_count,
                "duplicate_chunk_count": last_result.duplicate_chunk_count,
                "retrieval_verification": last_result.retrieval_verification,
                "vector_collection_name": self.vector_collection_name,
                "vector_chroma_dir": str(self.vector_chroma_dir),
                "steps": [asdict(step) for step in last_result.steps],
            }
            manifest["last_retrieval_verification"] = last_result.retrieval_verification
            manifest["pipeline"] = [asdict(step) for step in last_result.steps]
        else:
            manifest["pipeline"] = self._default_pipeline()
        self.manifest_path.parent.mkdir(parents=True, exist_ok=True)
        self.manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    def _default_pipeline(self) -> list[dict[str, Any]]:
        return [
            asdict(PipelineStep(STEP_PARSE, "pending", "等待上传资料", "")),
            asdict(PipelineStep(STEP_CHUNK, "pending", "等待解析文本", f"{self.chunk_token_size} token / overlap {self.chunk_overlap_tokens}")),
            asdict(PipelineStep(STEP_VECTOR, "pending", "等待写入向量索引", self.vector_collection_name)),
            asdict(PipelineStep(STEP_KG, "pending", "等待抽取实体关系", "")),
            asdict(PipelineStep(STEP_RETRIEVAL, "pending", "等待检索验证", "")),
        ]

    def _base_manifest(self) -> dict[str, Any]:
        return {
            "store_type": "project_runtime_knowledge",
            "base_dir": str(self.base_dir),
            "upload_dir": str(self.uploads_dir),
            "manifest_path": str(self.manifest_path),
            "rag_chunks_path": str(self.chunks_path),
            "kg_dir": str(self.kg_dir),
            "vector_enabled": self.vector_enabled,
            "vector_collection_name": self.vector_collection_name,
            "vector_chroma_dir": str(self.vector_chroma_dir),
            "chunk_token_size": self.chunk_token_size,
            "chunk_overlap_tokens": self.chunk_overlap_tokens,
            "min_chunk_tokens": self.min_chunk_tokens,
            "dedupe_key": "content_hash",
            "document_count": 0,
            "rag_chunk_count": 0,
            "vector_chunk_count": 0,
            "vector_ready": False,
            "kg_entity_count": 0,
            "kg_relation_count": 0,
            "structured_document_count": 0,
            "structured_block_count": 0,
            "markdown_document_count": 0,
            "last_retrieval_verification": {},
            "pipeline": self._default_pipeline(),
        }
